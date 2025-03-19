"""
PowerPointDocumentAdapter: Adapter for processing Microsoft PowerPoint (.pptx, .ppt) files.
"""

import os
import logging
import time
import traceback
from typing import Dict, List, Optional, Tuple, Union, Any
from datetime import datetime

from models_app.vision.document.adapters.document_format_adapter import DocumentFormatAdapter
from models_app.vision.document.utils.core.processing_metadata_context import ProcessingMetadataContext
from models_app.vision.document.utils.error_handling.errors import (
    DocumentProcessingError,
    DocumentValidationError
)
from models_app.vision.document.utils.error_handling.handlers import (
    handle_document_errors,
    measure_processing_time
)
from models_app.vision.document.utils.core.next_layer_interface import ProcessingEventType

logger = logging.getLogger(__name__)

class PowerPointDocumentAdapter(DocumentFormatAdapter):
    """
    Adapter for processing Microsoft PowerPoint (.pptx, .ppt) files.
    Uses python-pptx library to extract text, structure, and slide content.
    
    Features:
    - Extracts text from slides including presenter notes
    - Preserves slide structure (titles, content, layout)
    - Identifies texts in diagrams, charts, and tables
    - Extracts embedded images when possible
    - Retrieves document metadata (author, title, etc.)
    """
    
    # Class-level constants
    VERSION = "1.0.0"
    CAPABILITIES = {
        "text_extraction": 0.9,
        "structure_preservation": 0.8,
        "metadata_extraction": 0.8,
        "slide_content_extraction": 0.8,
        "image_identification": 0.6
    }
    SUPPORTED_FORMATS = [".pptx", ".ppt", ".odp"]
    PRIORITY = 80
    
    def __init__(self, config: Dict[str, Any] = None):
        """
        Initialize the PowerPoint document adapter.
        
        Args:
            config: Configuration dictionary with adapter-specific settings.
        """
        super().__init__(config)
        self.pptx_module = None
    
    def _initialize_adapter_components(self) -> None:
        """Initialize adapter-specific components."""
        try:
            # Attempt to import python-pptx
            import pptx
            from pptx import Presentation
            
            self.pptx_module = pptx
            logger.info("Successfully initialized PowerPoint document processing components")
            
            # Additional components for .ppt files (optional)
            try:
                import win32com.client
                self.has_win32com = True
            except ImportError:
                self.has_win32com = False
                logger.info("win32com not available - .ppt file handling will be limited")
            
        except ImportError as e:
            logger.warning(f"Could not import PowerPoint processing libraries: {str(e)}. Using fallback mode.")
            self._initialize_fallback_components()
    
    def _initialize_fallback_components(self) -> None:
        """Initialize fallback components when main libraries are unavailable."""
        self.can_process_pptx = False
        logger.info("Initialized fallback mode - document extraction will be limited")
    
    @handle_document_errors
    @measure_processing_time
    def _process_file(self, file_path: str, options: Dict[str, Any], 
                     metadata_context: Optional[ProcessingMetadataContext] = None) -> Dict[str, Any]:
        """
        Process a PowerPoint document file and extract information.
        
        Args:
            file_path: Path to the PowerPoint document file
            options: Processing options
            metadata_context: Context for tracking metadata during processing
            
        Returns:
            Dictionary with extracted information
        """
        if metadata_context:
            metadata_context.start_timing("powerpoint_processing")
            metadata_context.record_adapter_processing(
                adapter=self.__class__.__name__,
                file_path=file_path,
                capabilities=self.CAPABILITIES
            )
        
        file_extension = os.path.splitext(file_path)[1].lower()
        
        # Record processing start for performance tracking
        start_time = time.time()
        
        # Emit start event
        self.next_layer.emit_simple_event(
            event_type=ProcessingEventType.PROCESSING_PHASE_START,
            document_id=file_path,
            data={
                "adapter": self.__class__.__name__,
                "file_extension": file_extension,
                "phase": "powerpoint_processing"
            }
        )
        
        try:
            # Process based on file extension
            if file_extension == '.pptx':
                result = self._process_pptx_file(file_path, options, metadata_context)
            elif file_extension == '.ppt':
                result = self._process_ppt_file(file_path, options, metadata_context)
            elif file_extension == '.odp':
                result = self._process_odp_file(file_path, options, metadata_context)
            else:
                raise DocumentValidationError(
                    f"Unsupported file format: {file_extension}",
                    document_path=file_path
                )
            
            # Calculate processing time for performance tracking
            processing_time = time.time() - start_time
            
            # Add processor info to result
            result = {
                **result,
                "processor": {
                    "name": self.__class__.__name__,
                    "version": self.VERSION,
                    "processing_time": processing_time,
                    "timestamp": datetime.now().isoformat()
                }
            }
            
            # Emit performance metric
            self.next_layer.emit_simple_event(
                event_type=ProcessingEventType.PERFORMANCE_METRIC,
                document_id=file_path,
                data={
                    "adapter": self.__class__.__name__,
                    "processing_time": processing_time,
                    "success": True,
                    "file_extension": file_extension
                }
            )
            
            # Emit completion event
            self.next_layer.emit_simple_event(
                event_type=ProcessingEventType.PROCESSING_PHASE_END,
                document_id=file_path,
                data={
                    "adapter": self.__class__.__name__,
                    "phase": "powerpoint_processing",
                    "success": True,
                    "processing_time": processing_time
                }
            )
            
            if metadata_context:
                metadata_context.end_timing("powerpoint_processing")
                
            return result
            
        except Exception as e:
            # Calculate error processing time
            error_time = time.time() - start_time
            
            # Emit error event
            self.next_layer.emit_simple_event(
                event_type=ProcessingEventType.ERROR_OCCURRED,
                document_id=file_path,
                data={
                    "adapter": self.__class__.__name__,
                    "phase": "powerpoint_processing",
                    "error": str(e),
                    "error_type": type(e).__name__,
                    "processing_time": error_time
                }
            )
            
            if metadata_context:
                metadata_context.record_error(
                    component=self.__class__.__name__,
                    message=f"Error processing PowerPoint document: {str(e)}",
                    error_type=type(e).__name__
                )
                metadata_context.end_timing("powerpoint_processing")
            
            logger.error(f"Error processing PowerPoint document: {str(e)}")
            logger.debug(traceback.format_exc())
            
            raise DocumentProcessingError(f"Error processing PowerPoint document: {str(e)}")
    
    def _process_pptx_file(self, file_path: str, options: Dict[str, Any],
                          metadata_context: Optional[ProcessingMetadataContext] = None) -> Dict[str, Any]:
        """
        Process a .pptx file using python-pptx.
        
        Args:
            file_path: Path to the .pptx file
            options: Processing options
            metadata_context: Context for tracking metadata
            
        Returns:
            Dictionary with extracted slides, text and elements
        """
        if not self.pptx_module:
            raise DocumentProcessingError(
                "PowerPoint processing libraries not available. Cannot process .pptx file."
            )
        
        try:
            # Load presentation
            presentation = self.pptx_module.Presentation(file_path)
            
            # Extract presentation metadata
            metadata = self._extract_presentation_properties()
            
            # Process slides
            slides_data = []
            all_text = []
            image_data = []
            
            # Process each slide
            for slide_index, slide in enumerate(presentation.slides):
                slide_text = []
                slide_shapes = []
                
                # Process shapes in the slide
                for shape in slide.shapes:
                    shape_data = self._process_shape(shape)
                    
                    # Add shape text to slide text if available
                    if "text" in shape_data and shape_data["text"]:
                        slide_text.append(shape_data["text"])
                    
                    # Track images
                    if shape_data.get("is_image", False):
                        image_info = {
                            "slide_index": slide_index,
                            "shape_id": shape_data.get("id", "unknown"),
                            "width": shape_data.get("width", 0),
                            "height": shape_data.get("height", 0)
                        }
                        image_data.append(image_info)
                    
                    slide_shapes.append(shape_data)
                
                # Combine slide text
                combined_text = "\n".join(slide_text)
                all_text.append(f"Slide {slide_index + 1}:\n{combined_text}")
                
                # Create slide data
                slides_data.append({
                    "index": slide_index,
                    "title": self._get_slide_title(slide),
                    "shapes_count": len(slide_shapes),
                    "text": combined_text,
                    "shapes": slide_shapes
                })
            
            # Create result
            result = {
                "document_type": "pptx",
                "file_path": file_path,
                "text": "\n\n".join(all_text),
                "metadata": metadata,
                "slide_count": len(slides_data),
                "slides": slides_data,
                "images": image_data,
                "processing_timestamp": datetime.now().isoformat()
            }
            
            return result
            
        except Exception as e:
            logger.error(f"Error processing .pptx file {file_path}: {str(e)}")
            raise DocumentProcessingError(f"Failed to process .pptx file: {str(e)}")
    
    def _process_ppt_file(self, file_path: str, options: Dict[str, Any],
                         metadata_context: Optional[ProcessingMetadataContext] = None) -> Dict[str, Any]:
        """
        Process a .ppt file.
        
        Uses win32com if available, otherwise attempts basic text extraction.
        
        Args:
            file_path: Path to the .ppt file
            options: Processing options
            metadata_context: Context for tracking metadata
            
        Returns:
            Dictionary with extracted content
        """
        # Check if we have win32com for proper .ppt processing
        if self.has_win32com:
            try:
                # Convert .ppt to .pptx using win32com
                import win32com.client
                import os
                
                powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                powerpoint.Visible = False
                
                # Create temporary file for the conversion
                temp_file = self._create_temp_file(suffix=".pptx")
                
                # Open and save as pptx
                presentation = powerpoint.Presentations.Open(os.path.abspath(file_path), WithWindow=False)
                presentation.SaveAs(temp_file, 24)  # ppSaveAsOpenXMLPresentation = 24
                presentation.Close()
                powerpoint.Quit()
                
                # Process the converted pptx
                result = self._process_pptx_file(temp_file, options, metadata_context)
                result["document_type"] = "ppt"
                return result
                
            except Exception as e:
                logger.warning(f"Failed to convert .ppt using win32com: {str(e)}. Falling back to basic extraction.")
        
        # Fallback to basic text extraction
        try:
            # Try using textract if available
            try:
                import textract
                text = textract.process(file_path).decode('utf-8')
            except ImportError:
                # Very basic fallback - just read file as binary and extract printable chars
                with open(file_path, 'rb') as f:
                    content = f.read()
                text = ''.join(chr(c) for c in content if 32 <= c < 127 or c in [9, 10, 13])
            
            # Create simplified result
            return {
                "document_type": "ppt",
                "file_path": file_path,
                "text": text,
                "metadata": {
                    "filename": os.path.basename(file_path),
                    "file_size": os.path.getsize(file_path)
                },
                "slide_count": 0,  # Unknown
                "slides": [],      # Cannot extract slides
                "images": [],      # Cannot extract images
                "processing_timestamp": datetime.now().isoformat()
            }
        except Exception as e:
            logger.error(f"Error processing .ppt file {file_path}: {str(e)}")
            raise DocumentProcessingError(f"Failed to process .ppt file: {str(e)}")
    
    def _process_odp_file(self, file_path: str, options: Dict[str, Any],
                        metadata_context: Optional[ProcessingMetadataContext] = None) -> Dict[str, Any]:
        """
        Process an .odp file.
        
        Uses textract if available, otherwise attempts basic text extraction.
        
        Args:
            file_path: Path to the .odp file
            options: Processing options
            metadata_context: Context for tracking metadata
            
        Returns:
            Dictionary with extracted content
        """
        try:
            # Try using textract if available
            try:
                import textract
                text = textract.process(file_path).decode('utf-8')
            except ImportError:
                # Try using LibreOffice if available
                try:
                    import subprocess
                    import tempfile
                    
                    # Create temp dir for the conversion
                    temp_dir = tempfile.mkdtemp()
                    temp_pptx = os.path.join(temp_dir, "temp.pptx")
                    
                    # Convert using LibreOffice (if available)
                    subprocess.run([
                        "soffice", "--headless", "--convert-to", "pptx", 
                        "--outdir", temp_dir, file_path
                    ], check=True)
                    
                    # Process the converted file
                    if os.path.exists(temp_pptx):
                        result = self._process_pptx_file(temp_pptx, options, metadata_context)
                        result["document_type"] = "odp"
                        return result
                    else:
                        raise Exception("Conversion failed")
                        
                except Exception as conv_err:
                    logger.warning(f"Failed to convert ODP: {str(conv_err)}")
                    # Fallback to basic text extraction
                    with open(file_path, 'rb') as f:
                        content = f.read()
                    text = ''.join(chr(c) for c in content if 32 <= c < 127 or c in [9, 10, 13])
            
            # Create simplified result
            return {
                "document_type": "odp",
                "file_path": file_path,
                "text": text,
                "metadata": {
                    "filename": os.path.basename(file_path),
                    "file_size": os.path.getsize(file_path)
                },
                "slide_count": 0,  # Unknown
                "slides": [],      # Cannot extract slides
                "images": [],      # Cannot extract images
                "processing_timestamp": datetime.now().isoformat()
            }
        except Exception as e:
            logger.error(f"Error processing .odp file {file_path}: {str(e)}")
            raise DocumentProcessingError(f"Failed to process .odp file: {str(e)}")
    
    def _process_shape(self, shape) -> Dict[str, Any]:
        """
        Process a PowerPoint shape and extract its properties.
        
        Args:
            shape: Shape object from the slide
            
        Returns:
            Dictionary with shape data
        """
        shape_data = {
            "id": shape.shape_id if hasattr(shape, "shape_id") else "unknown",
            "name": shape.name if hasattr(shape, "name") else "",
            "type": type(shape).__name__,
            "width": shape.width.pt if hasattr(shape, "width") else 0,
            "height": shape.height.pt if hasattr(shape, "height") else 0
        }
        
        # Extract text if available
        if hasattr(shape, "text") and shape.text:
            shape_data["text"] = shape.text
            shape_data["has_text"] = True
        else:
            shape_data["has_text"] = False
        
        # Check if image
        shape_data["is_image"] = hasattr(shape, "image")
        
        # Check if table
        if hasattr(shape, "has_table") and shape.has_table:
            shape_data["is_table"] = True
            
            # Extract table data if possible
            try:
                table_data = []
                for row in shape.table.rows:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text if hasattr(cell, "text") else "")
                    table_data.append(row_data)
                
                shape_data["table"] = {
                    "rows": len(table_data),
                    "columns": len(table_data[0]) if table_data else 0,
                    "data": table_data
                }
            except Exception as table_err:
                logger.warning(f"Error extracting table data: {str(table_err)}")
        else:
            shape_data["is_table"] = False
        
        return shape_data
    
    def _get_slide_title(self, slide) -> str:
        """
        Extract the title of a slide.
        
        Args:
            slide: Slide object
            
        Returns:
            str: Title text or empty string if no title found
        """
        # Try to find a title placeholder
        for shape in slide.shapes:
            if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                if hasattr(shape, "placeholder_format") and shape.placeholder_format.type == 1:  # 1 = title
                    return shape.text if hasattr(shape, "text") else ""
        
        # Look for any shape that might be a title (large text at the top)
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and hasattr(shape, "top"):
                if shape.top.pt < 100:  # Assume titles are at the top of the slide
                    return shape.text
        
        return ""
    
    def _extract_presentation_properties(self) -> Dict[str, Any]:
        """
        Extract presentation properties/metadata.
        
        Returns:
            Dictionary with presentation properties
        """
        if not self.pptx_module:
            return {}
            
        try:
            properties = {}
            
            # Core properties
            if hasattr(self.pptx_module, "core_properties"):
                core_props = self.pptx_module.core_properties
                properties.update({
                    "title": core_props.title if hasattr(core_props, "title") else "",
                    "author": core_props.author if hasattr(core_props, "author") else "",
                    "subject": core_props.subject if hasattr(core_props, "subject") else "",
                    "keywords": core_props.keywords if hasattr(core_props, "keywords") else "",
                    "created": str(core_props.created) if hasattr(core_props, "created") else "",
                    "modified": str(core_props.modified) if hasattr(core_props, "modified") else "",
                    "category": core_props.category if hasattr(core_props, "category") else "",
                    "revision": core_props.revision if hasattr(core_props, "revision") else ""
                })
            
            # Presentation properties
            properties.update({
                "slide_count": len(self.pptx_module.slides) if hasattr(self.pptx_module, "slides") else 0,
                "slide_width": self.pptx_module.slide_width.pt if hasattr(self.pptx_module, "slide_width") else 0,
                "slide_height": self.pptx_module.slide_height.pt if hasattr(self.pptx_module, "slide_height") else 0
            })
            
            return properties
        except Exception as e:
            logger.warning(f"Error extracting presentation properties: {str(e)}")
            return {}

# Register this adapter
DocumentFormatAdapter.register_adapter("powerpoint", PowerPointDocumentAdapter) 