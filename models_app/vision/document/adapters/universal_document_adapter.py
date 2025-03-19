"""
UniversalDocumentAdapter: A single adapter that handles multiple document formats.
Provides a unified interface for processing different document types.
"""

import os
import logging
import json
import csv
import re
from typing import Dict, List, Any, Optional, Tuple, Union
from abc import ABC, abstractmethod
import tempfile
from datetime import datetime
import time

from models_app.vision.document.adapters.document_format_adapter import DocumentFormatAdapter
from models_app.vision.utils.image_processing.core import load_image, convert_to_array, save_image
from models_app.vision.utils.image_processing.enhancement import enhance_for_ocr
from models_app.vision.utils.image_processing.detection import detect_text_regions, detect_images
from error_handlers.models_app_errors.vision_errors import (
    handle_document_processing_error, 
    measure_processing_time,
    validate_document_path,
    DocumentValidationError,
    DocumentFormatError
)
from analytics_app.utils import monitor_document_performance
from models_app.vision.document.utils.plugin_system import register_adapter
from PyPDF2.errors import PdfReadError
from models_app.vision.ocr.ocr_model_selector import OCRModelSelector
from models_app.vision.document.utils.format_converter import FormatConverter
from models_app.vision.utils.testing.dummy_models import DummyModelFactory
from models_app.vision.utils.image_processing.adapter_preprocess import (
    enhance_image_for_ocr, detect_optimal_preprocessing
)
from models_app.vision.document.utils.core.processing_metadata_context import ProcessingMetadataContext
from PIL import Image
from models_app.vision.colpali.processor import ColPaliProcessor
from models_app.vision.document.factory.document_type_detector import DocumentTypeDetector
from models_app.vision.document.analyzer.document_analyzer import DocumentAnalyzer
from models_app.vision.document.document_base_adapter import DocumentBaseAdapter
from models_app.vision.document.utils.chunking_service import ChunkingService
from models_app.vision.document.utils.io.io_service import IOService
from models_app.vision.document.utils.error_handling.handlers import (
    handle_document_errors,
    handle_adapter_errors,
    measure_processing_time
)

logger = logging.getLogger(__name__)

# Try to import document processing libraries
try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    try:
        import pdfplumber
        PDF_AVAILABLE = True
    except ImportError:
        PDF_AVAILABLE = False


@register_adapter(name="universal", info={
    "description": "Universal Document Adapter for processing various document formats",
    "version": "1.0.0",
    "capabilities": {
        "office_documents": True,
        "pdfs": True,
        "text_files": True,
        "markup": True,
        "spreadsheets": True,
        "presentations": True
    },
    "priority": 80
})
class UniversalDocumentAdapter(DocumentBaseAdapter):
    """
    A unified adapter that handles multiple document formats through a common interface.
    Supports Word, Excel, PowerPoint, PDF, text files, and more.
    """
    
    def __init__(self, config: Dict[str, Any] = None):
        """
        Initialize the universal document adapter.
        
        Args:
            config: Configuration dictionary with adapter-specific settings.
        """
        super().__init__(config)
        self.format_name = "Universal Format"
        
        # Supported file extensions
        self._supported_extensions = [
            # Office documents
            ".docx", ".doc", 
            ".xlsx", ".xls",
            ".pptx", ".ppt",
            # PDF documents
            ".pdf",
            # Text documents
            ".txt", ".csv", ".json", ".md", ".html", ".xml",
            # Other formats
            ".rtf", ".odt", ".ods", ".odp"
        ]
        
        # Format handlers mapping (extension -> processing method)
        self._format_handlers = {
            # Office handlers
            ".docx": self._process_docx,
            ".doc": self._process_doc,
            ".xlsx": self._process_xlsx,
            ".xls": self._process_xls,
            ".pptx": self._process_pptx,
            ".ppt": self._process_ppt,
            # PDF handler
            ".pdf": self._process_pdf,
            # Text handlers
            ".txt": self._process_text,
            ".csv": self._process_csv,
            ".json": self._process_json,
            ".md": self._process_markdown,
            ".html": self._process_html,
            ".xml": self._process_xml,
            # Other handlers
            ".rtf": self._process_rtf,
            ".odt": self._process_odt,
            ".ods": self._process_ods,
            ".odp": self._process_odp
        }
        
        # Initialize handlers for each format
        self._init_format_handlers()
        
        # Initialisiere ChunkingService
        self.chunking_service = ChunkingService(self.config.get("chunking_config", None))
    
    def _init_format_handlers(self):
        """Initialize format handlers based on available libraries."""
        # Disable handlers if required libraries aren't available
        if not DOCX_AVAILABLE:
            self._format_handlers[".docx"] = self._process_unsupported
            self._format_handlers[".doc"] = self._process_unsupported
        
        if not EXCEL_AVAILABLE:
            self._format_handlers[".xlsx"] = self._process_unsupported
            self._format_handlers[".xls"] = self._process_unsupported
        
        if not PPTX_AVAILABLE:
            self._format_handlers[".pptx"] = self._process_unsupported
            self._format_handlers[".ppt"] = self._process_unsupported
        
        if not PDF_AVAILABLE:
            self._format_handlers[".pdf"] = self._process_unsupported
    
    def initialize(self) -> Dict[str, Any]:
        """
        Initialize the universal document adapter.
        
        Returns:
            Dict: Initialization status and metadata.
        """
        # Check which formats we can support based on available libraries
        available_formats = []
        
        if DOCX_AVAILABLE:
            available_formats.extend([".docx", ".doc"])
        
        if EXCEL_AVAILABLE:
            available_formats.extend([".xlsx", ".xls"])
        
        if PPTX_AVAILABLE:
            available_formats.extend([".pptx", ".ppt"])
        
        if PDF_AVAILABLE:
            available_formats.append(".pdf")
        
        # Text formats are always available
        available_formats.extend([".txt", ".csv", ".json", ".md", ".html", ".xml"])
        
        # Filter supported extensions to only include available ones
        self._supported_extensions = [ext for ext in self._supported_extensions 
                                     if ext in available_formats or ext in [".txt", ".csv", ".json"]]
        
        self.is_initialized = True
        return {
            "initialized": True,
            "metadata": {
                "supported_formats": self._supported_extensions,
                "libraries": {
                    "docx": DOCX_AVAILABLE,
                    "excel": EXCEL_AVAILABLE,
                    "pptx": PPTX_AVAILABLE,
                    "pdf": PDF_AVAILABLE
                }
            }
        }
    

    @monitor_document_performance
    @handle_document_processing_error
    def process_document(self, document_path: str, options: Dict[str, Any] = None,
                        metadata_context: Optional[ProcessingMetadataContext] = None) -> Dict[str, Any]:
        """
        Process a document and extract text, structure, and metadata.
        
        Args:
            document_path: Path to the document file.
            options: Additional options for processing.
            metadata_context: Context for tracking metadata during processing.
            
        Returns:
            Dict: Extracted content and metadata.
        """
        options = options or {}
        
        # Erstelle einen neuen Metadatenkontext, wenn keiner übergeben wurde
        if metadata_context is None:
            metadata_context = ProcessingMetadataContext(document_path)
            
        metadata_context.start_timing("universal_document_adapter_processing")
        
        try:
            if not self.is_initialized:
                metadata_context.start_timing("initialization")
                init_result = self.initialize()
                metadata_context.end_timing("initialization")
                
                if not init_result["initialized"]:
                    error_msg = "Failed to initialize universal document adapter"
                    metadata_context.record_error(
                        component="UniversalDocumentAdapter",
                        message=error_msg,
                        error_type="InitializationError",
                        is_fatal=True
                    )
                    metadata_context.end_timing("universal_document_adapter_processing")
                    return {
                        "error": error_msg,
                        "details": init_result.get("metadata", {})
                    }
            
            # Get file extension
            _, ext = os.path.splitext(document_path)
            ext = ext.lower()
            
            # Dokumenttyp in Metadaten aufzeichnen
            metadata_context.add_document_metadata("file_extension", ext)
            
            # Check if format is supported
            if ext not in self._supported_extensions:
                error_msg = f"Unsupported document format: {ext}"
                metadata_context.record_error(
                    component="UniversalDocumentAdapter",
                    message=error_msg,
                    error_type="UnsupportedFormat",
                    is_fatal=True
                )
                metadata_context.end_timing("universal_document_adapter_processing")
                return {
                    "error": error_msg,
                    "model": self._processor_name
                }
            
            # Get appropriate format handler
            handler = self._format_handlers.get(ext, self._process_unsupported)
            
            # Aufzeichnen, welcher Handler verwendet wird
            metadata_context.record_decision(
                component="UniversalDocumentAdapter",
                decision=f"Using handler for {ext}",
                reason=f"File extension matches {ext}"
            )
            
            # Process document using the handler
            metadata_context.start_timing(f"process_{ext}")
            result = handler(document_path, options, metadata_context)
            metadata_context.end_timing(f"process_{ext}")
            
            # Add standard metadata
            result.update({
                "model": self._processor_name,
                "format": self.format_name,
                "file_type": ext
            })
            
            # Extract knowledge graph if requested
            if options.get("extract_knowledge_graph", False):
                metadata_context.start_timing("knowledge_graph_extraction")
                knowledge_graph = self.extract_knowledge_graph(document_path, metadata_context)
                metadata_context.end_timing("knowledge_graph_extraction")
                result["knowledge_graph"] = knowledge_graph
            
            # Add metadata if not already present
            if "metadata" not in result:
                metadata_context.start_timing("metadata_extraction")
                result["metadata"] = self.extract_metadata(document_path, metadata_context)
                metadata_context.end_timing("metadata_extraction")
            
            metadata_context.end_timing("universal_document_adapter_processing")
            return result
            
        except Exception as e:
            error_msg = f"Error processing document: {str(e)}"
            logger.error(error_msg)
            metadata_context.record_error(
                component="UniversalDocumentAdapter",
                message=error_msg,
                error_type=type(e).__name__,
                is_fatal=True
            )
            metadata_context.end_timing("universal_document_adapter_processing")
            return {"error": error_msg}
    
    @handle_document_processing_error
    def extract_structure(self, document_path: str, options: Dict[str, Any] = None) -> Dict[str, Any]:
        """
        Extract the document's structure (headings, paragraphs, tables, etc.).
        
        Args:
            document_path: Path to the document file.
            options: Additional options for processing.
            
        Returns:
            Dict: Structured representation of the document.
        """
        options = options or {}
        
        try:
            # Get file extension
            _, ext = os.path.splitext(document_path)
            ext = ext.lower()
            
            # Get appropriate format handler
            handler = self._format_handlers.get(ext, self._process_unsupported)
            
            # Process document using the handler to get structure
            result = handler(document_path, {"extract_structure_only": True})
            
            # For detecting embedded images if document is an image or has images
            if options.get('preprocess_images', True):
                try:
                    document_image = load_image(document_path)
                    if document_image:
                        embedded_images = detect_images(document_image)
                        # We're not using embedded_images here, but it could be added to the result
                except Exception as img_error:
                    logger.warning(f"Could not process document as image: {str(img_error)}")
            
            return result.get("structure", {"error": "No structure extracted"})
            
        except Exception as e:
            logger.error(f"Error extracting structure from document {document_path}: {str(e)}")
            return {"error": f"Error extracting structure: {str(e)}"}
    
    def prepare_for_extraction(self, document_path: str, options: Dict[str, Any] = None,
                             metadata_context: Optional[ProcessingMetadataContext] = None, **kwargs) -> Dict[str, Any]:
        """
        Prepares a document for extraction using a multi-processor approach
        
        Args:
            document_path: Path to the document
            options: Dictionary of options for preparation
            metadata_context: Metadata context for tracking processing details
            **kwargs: Additional keyword arguments
            
        Returns:
            Dictionary with extraction preparation results
        """
        start_time = time.time()
        options = options or {}
        result = {
            "document_path": document_path,
            "prepared_at": datetime.now().isoformat(),
            "success": True
        }
        
        # Record adapter selection in context if available
        if metadata_context:
            metadata_context.start_timing("universal_adapter_preparation")
            metadata_context.record_adapter_selection(
                adapter_name="UniversalDocumentAdapter",
                reason="Selected for document preparation",
                confidence=0.7
            )
            
            # Record adapter capabilities
            capabilities = self.get_capabilities()
            metadata_context.record_capability_match(
                adapter_name="UniversalDocumentAdapter",
                capabilities=capabilities,
                match_score=0.7
            )
        
        try:
            # Check if we have document type and structure analysis in metadata_context
            document_type = None
            document_structure = None
            
            # Check options first (passed from DocumentProcessingManager)
            if options and "document_type" in options:
                document_type = options["document_type"]
                if metadata_context:
                    metadata_context.record_decision(
                        component="UniversalDocumentAdapter.prepare_for_extraction",
                        decision=f"Using document type from options: {document_type}",
                        reason="Document type passed from manager"
                    )
                    
            if options and "document_structure" in options:
                document_structure = options["document_structure"]
                if metadata_context:
                    metadata_context.record_decision(
                        component="UniversalDocumentAdapter.prepare_for_extraction",
                        decision="Using document structure from options",
                        reason="Structure passed from manager"
                    )
                    
            # If not in options, check metadata context
            if not document_type and metadata_context and metadata_context.has_analysis_result("document_type"):
                doc_type_info = metadata_context.get_analysis_result("document_type", {})
                document_type = doc_type_info.get("type")
                result["document_type_confidence"] = doc_type_info.get("confidence", 0.0)
                
                if metadata_context:
                    metadata_context.record_decision(
                        component="UniversalDocumentAdapter.prepare_for_extraction",
                        decision=f"Using document type from metadata context: {document_type}",
                        reason="Document type found in metadata context"
                    )
                    
            if not document_structure and metadata_context and metadata_context.has_analysis_result("document_structure"):
                document_structure = metadata_context.get_analysis_result("document_structure", {})
                
                if metadata_context:
                    metadata_context.record_decision(
                        component="UniversalDocumentAdapter.prepare_for_extraction",
                        decision="Using document structure from metadata context",
                        reason="Structure found in metadata context"
                    )
            
            # If we still don't have document type or structure, perform analysis
            if not document_type or not document_structure:
                if metadata_context:
                    metadata_context.start_timing("document_analysis")
                    metadata_context.record_decision(
                        component="UniversalDocumentAdapter.prepare_for_extraction",
                        decision="Performing document analysis",
                        reason="Missing document type or structure"
                    )
                
                # Only do document type detection if needed
                if not document_type:
                    detector = DocumentTypeDetector()
                    type_result = detector.classify_document_type(document_path)
                    document_type = type_result.get("type", "unknown")
                    result["document_type_confidence"] = type_result.get("confidence", 0.0)
                    
                    if metadata_context:
                        metadata_context.record_analysis_result("document_type", type_result)
                
                # Only do structure analysis if needed
                if not document_structure:
                    analyzer = DocumentAnalyzer()
                    analysis_result = analyzer.analyze_document(document_path)
                    document_structure = analysis_result.get("structure", {})
                    
                    if metadata_context:
                        metadata_context.record_analysis_result("document_structure", document_structure)
                        metadata_context.record_analysis_result("document_analysis", analysis_result)
                
                if metadata_context:
                    metadata_context.end_timing("document_analysis")
            
            # Set document type in result
            result["document_type"] = document_type
            
            # Select the appropriate processor based on document type
            if document_type in self.processors:
                processor = self.processors[document_type]
                if metadata_context:
                    metadata_context.record_decision(
                        component="UniversalDocumentAdapter.prepare_for_extraction",
                        decision=f"Using specialized processor for document type: {document_type}",
                        reason="Document type matched to registered processor"
                    )
            else:
                # Fall back to the default processor
                processor = self.processors["default"]
                if metadata_context:
                    metadata_context.record_decision(
                        component="UniversalDocumentAdapter.prepare_for_extraction",
                        decision="Using default processor",
                        reason=f"No specialized processor found for document type: {document_type}"
                    )
            
            # Start timing for processor preparation
            if metadata_context:
                metadata_context.start_timing("processor_preparation")
            
            # Pass document structure to processor options if available
            if document_structure:
                if "processor_options" not in options:
                    options["processor_options"] = {}
                options["processor_options"]["document_structure"] = document_structure
            
            # Call the processor's prepare method
            processor_result = processor.prepare(document_path, options)
            
            # Add processor result to our result
            result["processor_result"] = processor_result
            result["processed_file_path"] = processor_result.get("processed_file_path", document_path)
            
            # End timing for processor preparation
            if metadata_context:
                metadata_context.end_timing("processor_preparation")
                
                # Record processor decisions
                metadata_context.record_decision(
                    component="UniversalDocumentAdapter.prepare_for_extraction",
                    decision=f"Processor preparation complete: {processor.__class__.__name__}",
                    reason="Document successfully prepared by selected processor"
                )
            
            # Add processing time to result
            processing_time = time.time() - start_time
            result["processing_time"] = processing_time
            
            # Add metadata 
            result["metadata"] = {
                "filename": os.path.basename(document_path),
                "document_type": document_type,
                "preparation_time": processing_time,
                "processor_used": processor.__class__.__name__
            }
            
            return result
            
        except Exception as e:
            logger.error(f"Error preparing document for extraction: {str(e)}")
            if metadata_context:
                metadata_context.record_error(
                    component="UniversalDocumentAdapter",
                    message=f"Preparation for extraction failed: {str(e)}",
                    error_type=type(e).__name__
                )
            
            error_result = {
                "error": f"Preparation failed: {str(e)}",
                "document_path": document_path,
                "prepared_at": datetime.now().isoformat(),
                "success": False
            }
            
            return error_result
    
    def _extract_entities(self, text: str) -> List[Dict[str, Any]]:
        """Enhanced entity extraction with type classification"""
        entities = []
        
        # Use existing NER capabilities
        # ...existing code...
        
        # Add entity type classification
        for entity in basic_entities:
            entity_type = self._classify_entity_type(entity['text'], entity['context'])
            entity['type'] = entity_type
            entity['confidence'] = self._calculate_entity_confidence(entity)
            entities.append(entity)
        
        # Add domain-specific entity extraction
        domain_entities = self._extract_domain_entities(text)
        entities.extend(domain_entities)
        
        return entities
    
    def _classify_entity_type(self, entity_text: str, context: str) -> str:
        """Classify entity into detailed types (Person, Organization, Location, Date, etc.)"""
        # Implementation using rules or ML-based classification
        # ...
    
    def _extract_relationships(self, entities: List[Dict[str, Any]], 
                              structure: Dict[str, Any], 
                              metadata: Dict[str, Any]) -> List[Dict[str, Any]]:
        """
        Extract relationships between entities.
        
        Args:
            entities: List of extracted entities.
            structure: Document structure.
            metadata: Document metadata.
            
        Returns:
            List of relationships between entities.
        """
        relationships = []
        
        # Extract author relationship from metadata if available
        if "author" in metadata and entities:
            for entity in entities:
                if entity["type"] == "person" and metadata["author"] in entity["text"]:
                    relationships.append({
                        "source": entity["id"],
                        "target": "document",
                        "type": "author_of",
                        "confidence": 0.9
                    })
        
        # Look for entities that appear in the same paragraph
        paragraphs = []
        if "elements" in structure:
            paragraphs = [e for e in structure["elements"] if e.get("type") == "paragraph"]
        
        for paragraph in paragraphs:
            paragraph_entities = []
            paragraph_text = paragraph.get("text", "")
            
            # Find entities in this paragraph
            for entity in entities:
                if entity["text"] in paragraph_text:
                    paragraph_entities.append(entity)
            
            # Create relationships between entities in the same paragraph
            for i, entity1 in enumerate(paragraph_entities):
                for entity2 in paragraph_entities[i+1:]:
                    relationships.append({
                        "source": entity1["id"],
                        "target": entity2["id"],
                        "type": "appears_with",
                        "confidence": 0.7,
                        "context": "same_paragraph"
                    })
        
        return relationships
    
    def chunk_document(self, structure: Dict[str, Any], strategy: str = 'adaptive') -> List[Dict[str, Any]]:
        """
        Chunk the document structure into manageable pieces.
        
        Args:
            structure: Document structure to chunk
            strategy: Chunking strategy ('adaptive', 'fixed', or 'semantic').
            
        Returns:
            List of chunks with text and metadata
        """
        # Extract text content from structure
        text = structure.get('text', '')
        metadata = {k: v for k, v in structure.items() if k != 'text'}
        
        # Use the centralized ChunkingService instead of our own implementation
        return self.chunking_service.chunk_document(text, metadata, strategy)
    
    # The following methods have been moved to ChunkingService
    # They are kept here for reference but will be removed in the future
    """
    def _adaptive_chunking(self, structure: Dict[str, Any]) -> List[Dict[str, Any]]:
        # Implementation moved to ChunkingService
        pass
        
    def _fixed_chunking(self, structure: Dict[str, Any]) -> List[Dict[str, Any]]:
        # Implementation moved to ChunkingService
        pass
        
    def _semantic_chunking(self, structure: Dict[str, Any]) -> List[Dict[str, Any]]:
        # Implementation moved to ChunkingService
        pass
    """
    
    #
    # Format-specific handlers
    #
    
    def _process_docx(self, document_path: str, options: Dict[str, Any] = None, metadata_context: Optional[ProcessingMetadataContext] = None) -> Dict[str, Any]:
        """Process a Word document (.docx)."""
        options = options or {}
        
        if not DOCX_AVAILABLE:
            error_msg = "python-docx library not available. Install with pip install python-docx"
            if metadata_context:
                metadata_context.record_error(
                    component="UniversalDocumentAdapter._process_docx",
                    message=error_msg,
                    error_type="DependencyError"
                )
            return {"error": error_msg}
        
        try:
            if metadata_context:
                metadata_context.start_timing("docx_processing")
            
            # Open the document
            doc = docx.Document(document_path)
            
            # Extract text
            full_text = ""
            for para in doc.paragraphs:
                full_text += para.text + "\n"
            
            # Extract metadata
            metadata = {
                "title": doc.core_properties.title or "",
                "author": doc.core_properties.author or "",
                "created": str(doc.core_properties.created) if doc.core_properties.created else "",
                "modified": str(doc.core_properties.modified) if doc.core_properties.modified else "",
                "last_modified_by": doc.core_properties.last_modified_by or "",
                "revision": doc.core_properties.revision or 0,
                "page_count": len(doc.sections)
            }
            
            if metadata_context:
                metadata_context.add_document_metadata("docx_metadata", metadata)
            
            # Extract structure
            structure = {"elements": []}
            
            # Process paragraphs
            for i, para in enumerate(doc.paragraphs):
                if para.text.strip():
                    element = {
                        "type": "paragraph",
                        "text": para.text,
                        "index": i
                    }
                    
                    # Check for heading style
                    if para.style.name.startswith('Heading'):
                        element["type"] = "heading"
                        element["level"] = int(para.style.name.replace('Heading ', '')) if para.style.name != 'Heading' else 1
                    
                    structure["elements"].append(element)
            
            # Process tables
            tables = []
            for i, table in enumerate(doc.tables):
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                
                tables.append({
                    "index": i,
                    "data": table_data
                })
            
            if tables:
                structure["tables"] = tables
            
            # Extract images if requested
            image_paths = []
            if options.get("extract_images", False):
                if metadata_context:
                    metadata_context.start_timing("image_extraction")
                
                temp_dir = tempfile.mkdtemp()
                
                # Extract images from document
                image_index = 0
                for rel in doc.part.rels.values():
                    if "image" in rel.target_ref:
                        try:
                            image_data = rel.target_part.blob
                            image_path = os.path.join(temp_dir, f"image_{image_index}.png")
                            with open(image_path, "wb") as f:
                                f.write(image_data)
                            image_paths.append(image_path)
                            image_index += 1
                        except Exception as e:
                            logger.warning(f"Failed to extract image: {str(e)}")
                
                if metadata_context:
                    metadata_context.end_timing("image_extraction")
                    metadata_context.record_analysis_result("extracted_images_count", len(image_paths))
            
            # Process with OCR if requested
            if options.get("apply_ocr", False) and image_paths:
                if metadata_context:
                    metadata_context.start_timing("ocr_processing")
                
                ocr_text = self._apply_ocr_to_images(image_paths, options, metadata_context)
                
                if metadata_context:
                    metadata_context.end_timing("ocr_processing")
                
                if ocr_text:
                    # Combine extracted text and OCR text
                    full_text = f"{full_text}\n\n--- OCR Text ---\n\n{ocr_text}"
            
            if metadata_context:
                metadata_context.end_timing("docx_processing")
            
            return {
                "text": full_text,
                "structure": structure,
                "metadata": metadata,
                "images": image_paths
            }
            
        except Exception as e:
            error_msg = f"Error processing Word document: {str(e)}"
            logger.error(error_msg)
            
            if metadata_context:
                metadata_context.record_error(
                    component="UniversalDocumentAdapter._process_docx",
                    message=error_msg,
                    error_type=type(e).__name__
                )
                metadata_context.end_timing("docx_processing")
            
            return {
                "error": error_msg,
                "text": "",
                "metadata": {}
            }
    
    def _process_doc(self, document_path: str, options: Dict[str, Any] = None) -> Dict[str, Any]:
        """Process a legacy Word document (.doc)."""
        # This would typically use a converter to convert .doc to .docx first
        return {"error": "Legacy .doc format not directly supported. Convert to .docx first."}
    
    def _process_xlsx(self, document_path: str, options: Dict[str, Any] = None) -> Dict[str, Any]:
        """Process an Excel document (.xlsx)."""
        options = options or {}
        
        if not EXCEL_AVAILABLE:
            return {"error": "openpyxl library is not available"}
        
        try:
            workbook = openpyxl.load_workbook(document_path, data_only=True)
            
            # Extract structure (sheets and tables)
            sheets = []
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                
                # Extract rows as text
                rows = []
                for row in sheet.iter_rows():
                    row_data = []
                    for cell in row:
                        row_data.append(str(cell.value) if cell.value is not None else "")
                    rows.append(row_data)
                
                sheets.append({
                    "name": sheet_name,
                    "rows": rows,
                    "row_count": sheet.max_row,
                    "column_count": sheet.max_column
                })
            
            # Extract metadata
            metadata = {
                "title": workbook.properties.title or "",
                "author": workbook.properties.creator or "",
                "created": workbook.properties.created.isoformat() if workbook.properties.created else "",
                "modified": workbook.properties.modified.isoformat() if workbook.properties.modified else "",
                "last_modified_by": workbook.properties.lastModifiedBy or "",
                "category": workbook.properties.category or "",
                "keywords": workbook.properties.keywords or "",
                "subject": workbook.properties.subject or "",
                "content_type": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            }
            
            # Compile structure
            structure = {
                "sheets": sheets,
                "metadata": metadata
            }
            
            # If only structure is requested, return it
            if options.get("extract_structure_only", False):
                return {"structure": structure}
            
            # Extract text (simplified representation of tabular data)
            text = ""
            for sheet in sheets:
                text += f"Sheet: {sheet['name']}\n\n"
                for row in sheet['rows']:
                    text += " | ".join(row) + "\n"
                text += "\n\n"
            
            return {
                "text": text,
                "structure": structure,
                "metadata": metadata
            }
            
        except Exception as e:
            logger.error(f"Error processing Excel document {document_path}: {str(e)}")
            return {"error": f"Error processing Excel document: {str(e)}"}
    
    def _process_xls(self, document_path: str, options: Dict[str, Any] = None) -> Dict[str, Any]:
        """Process a legacy Excel document (.xls)."""
        # This would typically use a converter or library that supports older Excel formats
        return {"error": "Legacy .xls format not directly supported. Convert to .xlsx first."}
    
    def _process_pptx(self, document_path: str, options: Dict[str, Any] = None) -> Dict[str, Any]:
        """Process a PowerPoint document (.pptx)."""
        options = options or {}
        
        if not PPTX_AVAILABLE:
            return {"error": "python-pptx library is not available"}
        
        try:
            presentation = Presentation(document_path)
            
            # Extract structure (slides and shapes)
            slides = []
            for slide_idx, slide in enumerate(presentation.slides):
                # Extract shapes (text boxes, images, etc.)
                shapes = []
                for shape in slide.shapes:
                    shape_data = {
                        "name": shape.name,
                        "type": str(shape.shape_type)
                    }
                    
                    # Extract text if available
                    if hasattr(shape, "text"):
                        shape_data["text"] = shape.text
                    
                    shapes.append(shape_data)
                
                slides.append({
                    "index": slide_idx,
                    "shapes": shapes
                })
            
            # Extract metadata
            metadata = {
                "title": presentation.core_properties.title or "",
                "author": presentation.core_properties.author or "",
                "created": presentation.core_properties.created.isoformat() if presentation.core_properties.created else "",
                "modified": presentation.core_properties.modified.isoformat() if presentation.core_properties.modified else "",
                "last_modified_by": presentation.core_properties.last_modified_by or "",
                "revision": presentation.core_properties.revision or 0,
                "category": presentation.core_properties.category or "",
                "keywords": presentation.core_properties.keywords or "",
                "subject": presentation.core_properties.subject or "",
                "content_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            }
            
            # Compile structure
            structure = {
                "slides": slides,
                "metadata": metadata
            }
            
            # If only structure is requested, return it
            if options.get("extract_structure_only", False):
                return {"structure": structure}
            
            # Extract text from all slides
            text = ""
            for slide_idx, slide in enumerate(slides):
                text += f"Slide {slide_idx + 1}\n"
                for shape in slide["shapes"]:
                    if "text" in shape and shape["text"]:
                        text += shape["text"] + "\n"
                text += "\n"
            
            return {
                "text": text,
                "structure": structure,
                "metadata": metadata
            }
            
        except Exception as e:
            logger.error(f"Error processing PowerPoint document {document_path}: {str(e)}")
            return {"error": f"Error processing PowerPoint document: {str(e)}"}
    
    def _process_ppt(self, document_path: str, options: Dict[str, Any] = None) -> Dict[str, Any]:
        """Process a legacy PowerPoint document (.ppt)."""
        # This would typically use a converter to convert .ppt to .pptx first
        return {"error": "Legacy .ppt format not directly supported. Convert to .pptx first."}
    
    def _process_pdf(self, document_path: str, options: Dict[str, Any] = None) -> Dict[str, Any]:
        """
        Verarbeitet eine PDF-Datei und extrahiert Text und Metadaten.
        
        Args:
            document_path: Pfad zur PDF-Datei
            options: Verarbeitungsoptionen
        
        Returns:
            Dict: Extrahierte Inhalte und Metadaten
        """
        options = options or {}
        
        # 1. Extract images from PDF if needed
        extracted_images = []  # Initialize this variable to fix the undefined variable issue
        
        if options.get('extract_images', True):
            # Extract images from PDF to a temporary directory
            import tempfile
            temp_dir = tempfile.mkdtemp()
            
            try:
                # Use PDF extraction utility 
                from models_app.vision.document.utils.format_converter import FormatConverter
                converter = FormatConverter()
                extracted_images = converter.extract_images_from_pdf(document_path, output_dir=temp_dir)
            except Exception as e:
                logger.warning(f"Failed to extract images from PDF: {str(e)}")
        
        # Store extracted_images as image_paths for later use
        image_paths = extracted_images
        
        # 2. For image-based PDFs, enhance each extracted image
        if options.get('preprocess_images', True) and image_paths:
            for i, img_path in enumerate(image_paths):
                try:
                    _, np_image, _ = load_image(img_path)
                    enhanced_image = enhance_for_ocr(np_image)
                    save_image(enhanced_image, img_path)  # Replace with enhanced version
                except Exception as e:
                    logger.warning(f"Error enhancing PDF image {i}: {str(e)}")
        
        # Use PyMuPDF/fitz to extract text from PDF
        try:
            import fitz
        except ImportError:
            logger.error("PyMuPDF not installed. Install it with pip install pymupdf")
            return {
                "error": "PDF processing library not available",
                "text": "",
                "metadata": {}
            }
        
        try:
            # Open the PDF document
            doc = fitz.open(document_path)
            
            # Extract metadata
            metadata = {}
            for key, value in doc.metadata.items():
                if value:
                    metadata[key] = value
            
            # Extract text and structure from each page
            pages = []
            full_text = ""
            
            for page_idx, page in enumerate(doc):
                # Extract text from page
                page_text = page.get_text()
                
                # Extract images if available
                page_images = []
                if options.get('extract_images', True):
                    for img_idx, img in enumerate(page.get_images(full=True)):
                        xref = img[0]
                        base_image = doc.extract_image(xref)
                        if base_image:
                            image_info = {
                                "index": img_idx,
                                "width": base_image["width"],
                                "height": base_image["height"],
                                "referenced_in_page": page_idx
                            }
                            page_images.append(image_info)
                
                # Create page structure
                page_info = {
                    "index": page_idx,
                    "text": page_text,
                    "width": page.rect.width,
                    "height": page.rect.height,
                    "images": page_images
                }
                
                pages.append(page_info)
                full_text += page_text + "\n\n"
            
            # Create document structure
            structure = {
                "pages": pages,
                "metadata": metadata,
                "page_count": len(doc),
                "extracted_images": image_paths
            }
            
            # Close the document
            doc.close()
            
            # If only structure is requested, return it
            if options.get("extract_structure_only", False):
                return {"structure": structure}
            
            # Process with OCR if requested
            if options.get("apply_ocr", False) and image_paths:
                ocr_text = self._apply_ocr_to_images(image_paths, options, None)
                if ocr_text:
                    # Combine extracted text and OCR text
                    full_text = f"{full_text}\n\n--- OCR Text ---\n\n{ocr_text}"
            
            return {
                "text": full_text,
                "structure": structure,
                "metadata": metadata,
                "page_count": len(pages),
                "images": image_paths
            }
            
        except Exception as e:
            logger.error(f"Error processing PDF document {document_path}: {str(e)}")
            return {
                "error": f"Error processing PDF document: {str(e)}",
                "text": "",
                "metadata": {}
            }
    
    def _process_text(self, document_path: str, options: Dict[str, Any] = None) -> Dict[str, Any]:
        """Process a plain text document."""
        options = options or {}
        
        try:
            # Read text file
            with open(document_path, 'r', encoding='utf-8', errors='replace') as file:
                text = file.read()
            
            # Extract simple structure (paragraphs)
            paragraphs = [p for p in text.split('\n\n') if p.strip()]
            elements = [{"type": "paragraph", "text": p} for p in paragraphs]
            
            # Extract metadata
            metadata = self.extract_metadata(document_path)
            
            # Compile structure
            structure = {
                "elements": elements,
                "metadata": metadata
            }
            
            # If only structure is requested, return it
            if options.get("extract_structure_only", False):
                return {"structure": structure}
            
            return {
                "text": text,
                "structure": structure,
                "metadata": metadata
            }
            
        except Exception as e:
            logger.error(f"Error processing text document {document_path}: {str(e)}")
            return {"error": f"Error processing text document: {str(e)}"}
    
    def _process_csv(self, document_path: str, options: Dict[str, Any] = None) -> Dict[str, Any]:
        """Process a CSV document."""
        options = options or {}
        
        try:
            # Read CSV file
            rows = []
            with open(document_path, 'r', encoding='utf-8', errors='replace') as file:
                csv_reader = csv.reader(file)
                for row in csv_reader:
                    rows.append(row)
            
            # Extract header row if available
            header = rows[0] if rows else []
            data = rows[1:] if len(rows) > 1 else []
            
            # Compile structure
            structure = {
                "header": header,
                "data": data,
                "row_count": len(rows),
                "column_count": len(header) if header else 0
            }
            
            # Extract metadata
            metadata = self.extract_metadata(document_path)
            
            # If only structure is requested, return it
            if options.get("extract_structure_only", False):
                return {"structure": structure}
            
            # Convert to text
            text = ""
            if header:
                text += " | ".join(header) + "\n"
                text += "-" * (sum(len(h) + 3 for h in header)) + "\n"
            
            for row in data:
                text += " | ".join(row) + "\n"
            
            return {
                "text": text,
                "structure": structure,
                "metadata": metadata
            }
            
        except Exception as e:
            logger.error(f"Error processing CSV document {document_path}: {str(e)}")
            return {"error": f"Error processing CSV document: {str(e)}"}
    
    def _process_json(self, document_path: str, options: Dict[str, Any] = None) -> Dict[str, Any]:
        """Process a JSON document."""
        options = options or {}
        
        try:
            # Read JSON file
            with open(document_path, 'r', encoding='utf-8') as file:
                json_data = json.load(file)
            
            # Extract metadata
            metadata = self.extract_metadata(document_path)
            
            # Compile structure
            structure = {
                "data": json_data,
                "metadata": metadata
            }
            
            # If only structure is requested, return it
            if options.get("extract_structure_only", False):
                return {"structure": structure}
            
            # Convert to text (formatted JSON)
            text = json.dumps(json_data, indent=2)
            
            return {
                "text": text,
                "structure": structure,
                "metadata": metadata
            }
            
        except Exception as e:
            logger.error(f"Error processing JSON document {document_path}: {str(e)}")
            return {"error": f"Error processing JSON document: {str(e)}"}
    
    def _process_markdown(self, document_path: str, options: Dict[str, Any] = None) -> Dict[str, Any]:
        """Process a Markdown document."""
        options = options or {}
        
        try:
            # Read markdown file
            with open(document_path, 'r', encoding='utf-8', errors='replace') as file:
                text = file.read()
            
            # Simple structure extraction (headings, paragraphs)
            lines = text.split('\n')
            elements = []
            
            current_paragraph = []
            
            for line in lines:
                # Detect headings
                heading_match = re.match(r'^(#{1,6})\s+(.+)$', line)
                if heading_match:
                    # Save previous paragraph if it exists
                    if current_paragraph:
                        elements.append({
                            "type": "paragraph",
                            "text": '\n'.join(current_paragraph)
                        })
                        current_paragraph = []
                    
                    # Add heading
                    level = len(heading_match.group(1))
                    heading_text = heading_match.group(2)
                    elements.append({
                        "type": "heading",
                        "level": level,
                        "text": heading_text
                    })
                    
                # Empty line ends paragraph
                elif not line.strip() and current_paragraph:
                    elements.append({
                        "type": "paragraph",
                        "text": '\n'.join(current_paragraph)
                    })
                    current_paragraph = []
                    
                # Content line
                elif line.strip():
                    current_paragraph.append(line)
            
             # Add final paragraph if it exists
            if current_paragraph:
                elements.append({
                    "type": "paragraph",
                    "text": '\n'.join(current_paragraph)
                })
            
            # Extract metadata
            metadata = self.extract_metadata(document_path)
            
            # Compile structure
            structure = {
                "elements": elements,
                "metadata": metadata
            }
            
            # If only structure is requested, return it
            if options.get("extract_structure_only", False):
                return {"structure": structure}
            
            return {
                "text": text,
                "structure": structure,
                "metadata": metadata
            }
            
        except Exception as e:
            logger.error(f"Error processing Markdown document {document_path}: {str(e)}")
            return {"error": f"Error processing Markdown document: {str(e)}"}
    
    def _process_html(self, document_path: str, options: Dict[str, Any] = None) -> Dict[str, Any]:
        """Process an HTML document."""
        options = options or {}
        
        try:
            # Read HTML file
            with open(document_path, 'r', encoding='utf-8', errors='replace') as file:
                html_content = file.read()
            
            # Simple HTML parsing using regex (for production, use a proper parser like BeautifulSoup)
            # Extract title
            title_match = re.search(r'<title>(.*?)</title>', html_content, re.IGNORECASE | re.DOTALL)
            title = title_match.group(1) if title_match else ""
            
            # Extract body content
            body_match = re.search(r'<body>(.*?)</body>', html_content, re.IGNORECASE | re.DOTALL)
            body = body_match.group(1) if body_match else html_content
            
            # Remove HTML tags for text extraction
            text = re.sub(r'<[^>]*>', ' ', body)
            text = re.sub(r'\s+', ' ', text).strip()
            
            # Simple structure extraction
            elements = []
            
            # Extract headings
            for i in range(1, 7):
                heading_pattern = f'<h{i}>(.*?)</h{i}>'
                for match in re.finditer(heading_pattern, body, re.IGNORECASE | re.DOTALL):
                    heading_text = re.sub(r'<[^>]*>', '', match.group(1)).strip()
                    elements.append({
                        "type": "heading",
                        "level": i,
                        "text": heading_text
                    })
            
            # Extract paragraphs
            for match in re.finditer(r'<p>(.*?)</p>', body, re.IGNORECASE | re.DOTALL):
                paragraph_text = re.sub(r'<[^>]*>', '', match.group(1)).strip()
                if paragraph_text:
                    elements.append({
                        "type": "paragraph",
                        "text": paragraph_text
                    })
            
            # Extract metadata
            metadata = self.extract_metadata(document_path)
            metadata["title"] = title
            
            # Compile structure
            structure = {
                "elements": elements,
                "metadata": metadata
            }
            
            # If only structure is requested, return it
            if options.get("extract_structure_only", False):
                return {"structure": structure}
            
            return {
                "text": text,
                "structure": structure,
                "metadata": metadata
            }
            
        except Exception as e:
            logger.error(f"Error processing HTML document {document_path}: {str(e)}")
            return {"error": f"Error processing HTML document: {str(e)}"}
    
    def _process_xml(self, document_path: str, options: Dict[str, Any] = None) -> Dict[str, Any]:
        """Process an XML document."""
        options = options or {}
        
        try:
            # Read XML file
            with open(document_path, 'r', encoding='utf-8', errors='replace') as file:
                xml_content = file.read()
            
            # Simple extraction of text (for production, use a proper XML parser)
            # Remove XML tags for text extraction
            text = re.sub(r'<[^>]*>', ' ', xml_content)
            text = re.sub(r'\s+', ' ', text).strip()
            
            # Extract metadata
            metadata = self.extract_metadata(document_path)
            
            # Simple structure extraction
            elements = []
            
            # Extract XML elements (simplified)
            element_pattern = r'<([^\s/>]+)[^>]*>(.*?)</\1>'
            for match in re.finditer(element_pattern, xml_content, re.DOTALL):
                tag_name = match.group(1)
                content = match.group(2)
                
                # Remove nested tags for text extraction
                element_text = re.sub(r'<[^>]*>', '', content).strip()
                if element_text:
                    elements.append({
                        "type": "element",
                        "tag": tag_name,
                        "text": element_text
                    })
            
            # Compile structure
            structure = {
                "elements": elements,
                "metadata": metadata
            }
            
            # If only structure is requested, return it
            if options.get("extract_structure_only", False):
                return {"structure": structure}
            
            return {
                "text": text,
                "structure": structure,
                "metadata": metadata
            }
            
        except Exception as e:
            logger.error(f"Error processing XML document {document_path}: {str(e)}")
            return {"error": f"Error processing XML document: {str(e)}"}
    
    def _process_rtf(self, document_path: str, options: Dict[str, Any] = None) -> Dict[str, Any]:
        """Process an RTF document."""
        # This would typically use a RTF-to-text converter
        return {"error": "RTF format not directly supported. Convert to .docx or .txt first."}
    
    def _process_odt(self, document_path: str, options: Dict[str, Any] = None) -> Dict[str, Any]:
        """Process an OpenDocument Text document."""
        # This would typically use a library for OpenDocument formats
        return {"error": "OpenDocument Text format not directly supported. Convert to .docx first."}
    
    def _process_ods(self, document_path: str, options: Dict[str, Any] = None) -> Dict[str, Any]:
        """Process an OpenDocument Spreadsheet document."""
        # This would typically use a library for OpenDocument formats
        return {"error": "OpenDocument Spreadsheet format not directly supported. Convert to .xlsx first."}
    
    def _process_odp(self, document_path: str, options: Dict[str, Any] = None) -> Dict[str, Any]:
        """Process an OpenDocument Presentation document."""
        # This would typically use a library for OpenDocument formats
        return {"error": "OpenDocument Presentation format not directly supported. Convert to .pptx first."}
    
    def _process_unsupported(self, document_path: str, options: Dict[str, Any] = None) -> Dict[str, Any]:
        """Fallback for unsupported formats."""
        _, ext = os.path.splitext(document_path)
        return {
            "error": f"Unsupported document format: {ext}",
            "metadata": self.extract_metadata(document_path)
        }
    
    def get_capabilities(self) -> Dict[str, float]:
        """
        Returns capabilities of this adapter with confidence scores between 0.0 and 1.0.
        
        The capabilities dictionary maps capability names to confidence scores, where:
        - 0.0 means the adapter cannot handle this capability
        - 1.0 means the adapter is fully confident it can handle this capability
        - Values between represent partial confidence
        
        Returns:
            Dict[str, float]: Mapping of capability names to confidence scores
        """
        # Base capabilities for the universal adapter
        capabilities = {
            # Document types
            "text_extraction": 0.7,
            "office_documents": 0.8 if all([DOCX_AVAILABLE, EXCEL_AVAILABLE, PPTX_AVAILABLE]) else 0.3,
            "pdfs": 0.7 if PDF_AVAILABLE else 0.0,
            "images": 0.1,  # Low confidence for image processing
            "scanned_documents": 0.5 if PDF_AVAILABLE else 0.0,
            "text_files": 1.0,  # Highest confidence for basic text files
            
            # Content types
            "tables": 0.7 if EXCEL_AVAILABLE else 0.3,
            "structured_data": 0.8,
            "forms": 0.5,
            "diagrams": 0.2,
            "complex_layouts": 0.5 if PDF_AVAILABLE else 0.2,
            "mixed_content": 0.6,
            
            # Format specific
            "spreadsheets": 0.9 if EXCEL_AVAILABLE else 0.0,
            "presentations": 0.9 if PPTX_AVAILABLE else 0.0,
            "word_documents": 0.9 if DOCX_AVAILABLE else 0.0,
        }
        
        # Check which libraries are available and adjust scores
        if not self.is_initialized:
            # Lower all scores if not initialized
            capabilities = {k: max(v * 0.5, 0.1) for k, v in capabilities.items()}
        
        # Check configuration for special handling options
        if self.config.get("optimize_for_tables", False):
            capabilities["tables"] = min(capabilities["tables"] + 0.2, 1.0)
            capabilities["spreadsheets"] = min(capabilities["spreadsheets"] + 0.1, 1.0)
            
        if self.config.get("optimize_for_text", False):
            capabilities["text_extraction"] = min(capabilities["text_extraction"] + 0.2, 1.0)
            capabilities["text_files"] = 1.0
            
        return capabilities

    def _apply_ocr_to_images(self, image_paths, options=None, metadata_context=None):
        """
        Intelligente Bildverarbeitung mit ColPali als primärem Prozessor.
        
        Diese Methode verwendet ColPali als zentralen Prozessor für alle Bildtypen, da er sowohl
        Text in Bildern erkennen als auch Bildinhalte beschreiben kann. Für Bilder mit erkanntem
        Text wird zusätzlich ein spezialisierter OCR-Prozessor verwendet, um optimale Ergebnisse
        zu erzielen.
        
        Der UniversalDocumentAdapter ist primär für Office-Dateien und textbasierte Dokumente
        konzipiert, kann aber durch diese Methode auch eingebettete Bilder intelligent verarbeiten.
        Bei komplexen Dokumenten mit vielen visuellen Elementen sollte jedoch der HybridDocumentAdapter
        bevorzugt werden.
        
        Args:
            image_paths: Liste der Bildpfade
            options: Optionale Konfigurationsparameter
            metadata_context: Kontext für das Tracking von Metadaten während der Verarbeitung
            
        Returns:
            String mit kombinierten Ergebnissen aller Bilder
        """
        options = options or {}
        combined_results = []
        
        from models_app.vision.processors.colpali.processor import ColPaliProcessor
        from models_app.vision.ocr.ocr_model_selector import OCRModelSelector
        from models_app.vision.utils.image_processing.adapter_preprocess import (
            enhance_image_for_ocr, detect_optimal_preprocessing
        )
        import tempfile
        import os
        
        colpali = ColPaliProcessor()
        selector = OCRModelSelector()
        
        if metadata_context:
            metadata_context.start_timing("ocr_processing")
            metadata_context.record_decision(
                component="UniversalDocumentAdapter",
                decision="Using OCR for images",
                reason=f"Processing {len(image_paths)} images with OCR"
            )
        
        for img_path in image_paths:
            if metadata_context:
                metadata_context.start_timing(f"process_image_{os.path.basename(img_path)}")
            
            try:
                # Bestimme die optimale Vorverarbeitungsmethode für das Bild
                if metadata_context:
                    metadata_context.start_timing("preprocessing_detection")
                
                preprocessing_method = detect_optimal_preprocessing(img_path)
                
                if metadata_context:
                    metadata_context.end_timing("preprocessing_detection")
                    metadata_context.record_preprocessing_step(
                        method="detect_optimal_preprocessing",
                        parameters={"selected_method": preprocessing_method},
                        before_image_path=img_path
                    )
                
                # Wende die Vorverarbeitung an
                if metadata_context:
                    metadata_context.start_timing("image_enhancement")
                
                # Erstelle eine temporäre Datei für das verarbeitete Bild
                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as temp_file:
                    temp_path = temp_file.name
                
                # Verbessere das Bild für OCR
                enhanced_image = enhance_image_for_ocr(img_path, method=preprocessing_method)
                enhanced_image.save(temp_path)
                
                if metadata_context:
                    metadata_context.end_timing("image_enhancement")
                    metadata_context.record_preprocessing_step(
                        method="enhance_image_for_ocr",
                        parameters={"method": preprocessing_method},
                        before_image_path=img_path,
                        after_image_path=temp_path
                    )
                
                # Verwende ColPali für die primäre Analyse aller Bilder
                if metadata_context:
                    metadata_context.start_timing("colpali_processing")
                
                colpali_result = colpali.process_image(temp_path, options)
                
                if metadata_context:
                    metadata_context.end_timing("colpali_processing")
                    metadata_context.add_adapter_data(
                        adapter_name="UniversalDocumentAdapter",
                        key=f"colpali_result_{os.path.basename(img_path)}",
                        value={"contains_text": colpali_result.get("contains_text", False)}
                    )
                
                # ColPali kann selbst erkennen, ob Text im Bild vorhanden ist
                if colpali_result.get("contains_text", False):
                    # Bei erkanntem Text zusätzlich spezialisierten OCR-Prozessor verwenden
                    if metadata_context:
                        metadata_context.start_timing("ocr_selection")
                    
                    ocr_adapter = selector.select_model(temp_path)
                    
                    if metadata_context:
                        metadata_context.end_timing("ocr_selection")
                        metadata_context.record_ocr_selection(
                            engine_name=ocr_adapter.__class__.__name__,
                            reason="Selected based on image analysis"
                        )
                        metadata_context.start_timing("ocr_processing")
                    
                    # Übergebe das bereits vorverarbeitete Bild an den OCR-Adapter
                    ocr_options = options.copy()
                    ocr_options["already_preprocessed"] = True
                    ocr_result = ocr_adapter.process_image(temp_path, ocr_options)
                    
                    if metadata_context:
                        metadata_context.end_timing("ocr_processing")
                    
                    # Kombiniere OCR-Text mit ColPali-Kontext für bessere Ergebnisse
                    if ocr_result.get("text", "").strip():
                        combined_results.append(ocr_result.get("text", ""))
                    else:
                        # Fallback auf ColPali-Text, falls OCR nichts findet
                        combined_results.append(colpali_result.get("text", colpali_result.get("caption", "")))
                else:
                    # Für Bilder ohne Text verwende die ColPali-Bildbeschreibung
                    combined_results.append(colpali_result.get("caption", ""))
                
                # Lösche die temporäre Datei
                try:
                    os.unlink(temp_path)
                except Exception as e:
                    logger.warning(f"Failed to delete temporary file {temp_path}: {str(e)}")
                
            except Exception as e:
                error_msg = f"Error processing image {img_path}: {str(e)}"
                logger.error(error_msg)
                if metadata_context:
                    metadata_context.record_error(
                        component="UniversalDocumentAdapter._apply_ocr_to_images",
                        message=error_msg,
                        error_type=type(e).__name__
                    )
            
            if metadata_context:
                metadata_context.end_timing(f"process_image_{os.path.basename(img_path)}")
        
        if metadata_context:
            metadata_context.end_timing("ocr_processing")
        
        return "\n\n".join(combined_results)

    def extract_knowledge_graph(self, document_path, metadata_context=None):
        """Deprecated method, use prepare_for_extraction instead."""
        logger.warning("extract_knowledge_graph is deprecated, use prepare_for_extraction instead")
        return self.prepare_for_extraction(document_path, metadata_context)

    def _initialize_format_handlers(self):
        """Initialisiert Dummy-Format-Handler für verschiedene Dokumenttypen."""
        if self.config.get('dummy_mode', False):
            dummy = DummyModelFactory.create_document_dummy("universal")
            self.format_handlers = dummy.get("format_handlers")
        else:
            # Originaler Code zur Format-Handler-Initialisierung
            pass